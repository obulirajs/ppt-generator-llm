# Write Python script that can generate 3 PPTX templates based on given specification
script_content = '''
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

templates = [
    {
        'id': 1,
        'name': 'Corporate Default',
        'colors': {'accent': 'FF6600', 'primary': '003366', 'secondary': '0066CC'},
        'fonts': {'title': 'Arial Bold', 'title_size': 32, 'body': 'Calibri', 'body_size': 18},
        'description': 'Standard corporate template with company branding.'
    },
    {
        'id': 2,
        'name': 'Modern Pitch',
        'colors': {'accent': 'E94560', 'primary': '1A1A2E', 'secondary': '16213E'},
        'fonts': {'title': 'Montserrat', 'title_size': 36, 'body': 'Open Sans', 'body_size': 20},
        'description': 'Modern template for pitch presentations.'
    },
    {
        'id': 3,
        'name': 'Professional Report',
        'colors': {'accent': '3498DB', 'primary': '2C3E50', 'secondary': '34495E'},
        'fonts': {'title': 'Georgia', 'title_size': 30, 'body': 'Times New Roman', 'body_size': 16},
        'description': 'Professional template for business reports.'
    }
]

for t in templates:
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = t['name']
    subtitle.text = t['description']

    title.text_frame.paragraphs[0].font.name = t['fonts']['title']
    title.text_frame.paragraphs[0].font.size = Pt(t['fonts']['title_size'])
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(t['colors']['primary'])

    body_slide = prs.slides.add_slide(prs.slide_layouts[1])
    body_title = body_slide.shapes.title
    body_content = body_slide.placeholders[1]
    body_title.text = 'Sample Content'
    body_content.text = f"Accent Color: #{t['colors']['accent']}\nPrimary Color: #{t['colors']['primary']}\nSecondary Color: #{t['colors']['secondary']}"

    body_content.text_frame.paragraphs[0].font.name = t['fonts']['body']
    body_content.text_frame.paragraphs[0].font.size = Pt(t['fonts']['body_size'])

    file_name = f"template_{t['id']}_{t['name'].replace(' ', '_')}.pptx"
    prs.save(file_name)
    print(f"Saved: {file_name}")
'''

with open('generate_ppt_templates.py', 'w') as f:
    f.write(script_content)

'import os; os.listdir()'