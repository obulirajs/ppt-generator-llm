import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR

class PPTXTemplateGenerator:
    """Generate PowerPoint files based on template specifications."""
    
    def __init__(self, template_data):
        """Initialize with template data from JSON."""
        self.template_data = template_data
        
    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def create_presentation(self, template):
        """Create a PowerPoint presentation based on template specifications."""
        prs = Presentation()
        
        # Set slide size (16:9 widescreen)
        # Default PowerPoint slide size is already 16:9, no need to modify
        
        # Extract template settings
        colors = template['brand_colors']
        fonts = template['font_settings']
        
        # Create different slide layouts based on category
        if template['category'] == 'corporate':
            self.create_corporate_slides(prs, template)
        elif template['category'] == 'pitch':
            self.create_pitch_slides(prs, template)
        elif template['category'] == 'report':
            self.create_report_slides(prs, template)
        
        return prs
    
    def apply_text_formatting(self, text_frame, title_text, font_name, font_size, color_hex):
        """Apply formatting to a text frame."""
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = title_text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = self.hex_to_rgb(color_hex)
        return p
    
    def create_corporate_slides(self, prs, template):
        """Create slides for corporate template."""
        colors = template['brand_colors']
        fonts = template['font_settings']
        
        # Slide 1: Title Slide
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.hex_to_rgb(colors['primary'])
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        self.apply_text_formatting(
            title_frame, 
            f"{template['name']} Template",
            fonts['title'],
            fonts['size_title'],
            '#FFFFFF'
        )
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        self.apply_text_formatting(
            subtitle_frame,
            template['description'],
            fonts['body'],
            fonts['size_body'],
            colors['secondary']
        )
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Slide 2: Content Slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Header with accent color
        header = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(10), Inches(1)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.hex_to_rgb(colors['accent'])
        header.line.fill.background()
        
        # Header text
        header_text = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
        )
        self.apply_text_formatting(
            header_text.text_frame,
            "Section Title",
            fonts['title'],
            24,
            '#FFFFFF'
        )
        
        # Content area
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(3.5)
        )
        content_frame = content_box.text_frame
        p = self.apply_text_formatting(
            content_frame,
            "• Key Point 1: Professional corporate design",
            fonts['body'],
            fonts['size_body'],
            colors['primary']
        )
        
        # Add more bullet points
        p = content_frame.add_paragraph()
        p.text = "• Key Point 2: Consistent brand colors"
        p.font.name = fonts['body']
        p.font.size = Pt(fonts['size_body'])
        p.font.color.rgb = self.hex_to_rgb(colors['primary'])
        
        p = content_frame.add_paragraph()
        p.text = "• Key Point 3: Clean and modern layout"
        p.font.name = fonts['body']
        p.font.size = Pt(fonts['size_body'])
        p.font.color.rgb = self.hex_to_rgb(colors['primary'])
        
    def create_pitch_slides(self, prs, template):
        """Create slides for pitch template."""
        colors = template['brand_colors']
        fonts = template['font_settings']
        
        # Slide 1: Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Gradient-like effect with shapes
        bg_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(10), Inches(5.625)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.hex_to_rgb(colors['primary'])
        bg_shape.line.fill.background()
        
        # Accent shape
        accent_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(3.5), Inches(10), Inches(2.125)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = self.hex_to_rgb(colors['secondary'])
        accent_shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(2)
        )
        title_frame = title_box.text_frame
        self.apply_text_formatting(
            title_frame,
            "MODERN PITCH DECK",
            fonts['title'],
            fonts['size_title'],
            colors['accent']
        )
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Tagline
        tagline_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(8), Inches(1)
        )
        self.apply_text_formatting(
            tagline_box.text_frame,
            "Compelling • Innovative • Results-Driven",
            fonts['body'],
            fonts['size_body'] - 2,
            '#FFFFFF'
        )
        tagline_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Slide 2: Problem/Solution
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Split layout
        left_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(5), Inches(5.625)
        )
        left_shape.fill.solid()
        left_shape.fill.fore_color.rgb = self.hex_to_rgb(colors['accent'])
        left_shape.line.fill.background()
        
        # Problem text
        problem_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1), Inches(4), Inches(3.5)
        )
        problem_frame = problem_box.text_frame
        self.apply_text_formatting(
            problem_frame,
            "THE PROBLEM",
            fonts['title'],
            28,
            '#FFFFFF'
        )
        
        p = problem_frame.add_paragraph()
        p.text = "\nMarket challenges and pain points that need innovative solutions"
        p.font.name = fonts['body']
        p.font.size = Pt(fonts['size_body'])
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Solution text
        solution_box = slide.shapes.add_textbox(
            Inches(5.5), Inches(1), Inches(4), Inches(3.5)
        )
        solution_frame = solution_box.text_frame
        self.apply_text_formatting(
            solution_frame,
            "OUR SOLUTION",
            fonts['title'],
            28,
            colors['primary']
        )
        
        p = solution_frame.add_paragraph()
        p.text = "\nInnovative approach that addresses key market needs effectively"
        p.font.name = fonts['body']
        p.font.size = Pt(fonts['size_body'])
        p.font.color.rgb = self.hex_to_rgb(colors['secondary'])
        
    def create_report_slides(self, prs, template):
        """Create slides for report template."""
        colors = template['brand_colors']
        fonts = template['font_settings']
        
        # Slide 1: Cover Page
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Professional border
        border = slide.shapes.add_shape(
            1, Inches(0.25), Inches(0.25), Inches(9.5), Inches(5.125)
        )
        border.fill.background()
        border.line.color.rgb = self.hex_to_rgb(colors['primary'])
        border.line.width = 14400
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(1)
        )
        self.apply_text_formatting(
            title_box.text_frame,
            "PROFESSIONAL REPORT",
            fonts['title'],
            fonts['size_title'],
            colors['primary']
        )
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(0.5)
        )
        self.apply_text_formatting(
            subtitle_box.text_frame,
            template['description'],
            fonts['body'],
            fonts['size_body'],
            colors['secondary']
        )
        subtitle_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(8), Inches(0.5)
        )
        self.apply_text_formatting(
            date_box.text_frame,
            datetime.now().strftime("%B %Y"),
            fonts['body'],
            fonts['size_body'] - 2,
            colors['accent']
        )
        date_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Slide 2: Executive Summary
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Header
        header_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        self.apply_text_formatting(
            header_box.text_frame,
            "Executive Summary",
            fonts['title'],
            fonts['size_title'] - 4,
            colors['primary']
        )
        
        # Divider line
        line = slide.shapes.add_connector(
            1, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3)
        )
        line.line.color.rgb = self.hex_to_rgb(colors['accent'])
        line.line.width = 9525
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(3.5)
        )
        content_frame = content_box.text_frame
        
        # Key findings
        p = self.apply_text_formatting(
            content_frame,
            "Key Findings:",
            fonts['title'],
            fonts['size_body'] + 2,
            colors['secondary']
        )
        
        findings = [
            "Strategic analysis reveals significant opportunities",
            "Market conditions favor proposed approach",
            "Implementation timeline aligns with objectives",
            "Risk factors have been identified and mitigated"
        ]
        
        for finding in findings:
            p = content_frame.add_paragraph()
            p.text = f"  • {finding}"
            p.font.name = fonts['body']
            p.font.size = Pt(fonts['size_body'])
            p.font.color.rgb = self.hex_to_rgb(colors['primary'])
            p.space_after = Pt(6)
    
    def generate_all_templates(self, output_dir="."):
        """Generate PPTX files for all templates."""
        templates = self.template_data.get('templates', [])
        generated_files = []
        
        for template in templates:
            try:
                # Create presentation
                prs = self.create_presentation(template)
                
                # Generate filename
                filename = f"{output_dir}/{template['name'].replace(' ', '_').lower()}_template.pptx"
                
                # Save presentation
                prs.save(filename)
                generated_files.append({
                    'template_name': template['name'],
                    'filename': filename,
                    'category': template['category']
                })
                
                print(f"✓ Generated: {filename}")
                
            except Exception as e:
                print(f"✗ Error generating {template['name']}: {str(e)}")
        
        return generated_files


def main():
    """Main function to generate PowerPoint templates."""
    
    # JSON data from the uploaded file
    json_data = {
        "default_template_id": 1,
        "success": True,
        "templates": [
            {
                "brand_colors": {
                    "accent": "#FF6600",
                    "primary": "#003366",
                    "secondary": "#0066CC"
                },
                "category": "corporate",
                "created_at": "2025-10-23T14:33:12",
                "description": "Standard corporate template with company branding",
                "file_path": "placeholder.pptx",
                "font_settings": {
                    "body": "Calibri",
                    "size_body": 18,
                    "size_title": 32,
                    "title": "Arial"
                },
                "id": 1,
                "is_default": True,
                "is_mine": False,
                "is_public": False,
                "name": "Corporate Default",
                "thumbnail_path": None,
                "type": "system",
                "usage_count": 4
            },
            {
                "brand_colors": {
                    "accent": "#E94560",
                    "primary": "#1A1A2E",
                    "secondary": "#16213E"
                },
                "category": "pitch",
                "created_at": "2025-10-23T14:33:12",
                "description": "Modern template for pitch presentations",
                "file_path": "placeholder.pptx",
                "font_settings": {
                    "body": "Open Sans",
                    "size_body": 20,
                    "size_title": 36,
                    "title": "Montserrat"
                },
                "id": 2,
                "is_default": False,
                "is_mine": False,
                "is_public": False,
                "name": "Modern Pitch",
                "thumbnail_path": None,
                "type": "system",
                "usage_count": 1
            },
            {
                "brand_colors": {
                    "accent": "#3498DB",
                    "primary": "#2C3E50",
                    "secondary": "#34495E"
                },
                "category": "report",
                "created_at": "2025-10-23T14:33:12",
                "description": "Professional template for business reports",
                "file_path": "placeholder.pptx",
                "font_settings": {
                    "body": "Times New Roman",
                    "size_body": 16,
                    "size_title": 30,
                    "title": "Georgia"
                },
                "id": 3,
                "is_default": False,
                "is_mine": False,
                "is_public": False,
                "name": "Professional Report",
                "thumbnail_path": None,
                "type": "system",
                "usage_count": 0
            }
        ],
        "timestamp": "2025-10-25T14:05:55.714432",
        "total": 3
    }
    
    # Create generator instance
    generator = PPTXTemplateGenerator(json_data)
    
    # Generate all templates
    print("Starting PowerPoint template generation...\n")
    generated = generator.generate_all_templates()
    
    print(f"\n✅ Successfully generated {len(generated)} PowerPoint templates!")
    print("\nGenerated files:")
    for file_info in generated:
        print(f"  - {file_info['template_name']} ({file_info['category']}): {file_info['filename']}")


if __name__ == "__main__":
    main()