"""
PowerPoint Generator Module
Creates professional presentations from structured JSON data
Uses python-pptx library to generate .pptx files
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import os
import uuid
from datetime import datetime
import shutil

# ==================== CONFIGURATION ====================

# Output Configuration
OUTPUT_FOLDER = "temp"
FILENAME_PREFIX = "presentation"

# Color Scheme (RGB values)
TITLE_COLOR = RGBColor(31, 78, 121)           # Dark Blue
SUBTITLE_COLOR = RGBColor(89, 89, 89)         # Gray
CONTENT_TITLE_COLOR = RGBColor(44, 62, 80)    # Dark Slate
BULLET_COLOR = RGBColor(68, 68, 68)           # Dark Gray
BACKGROUND_COLOR = RGBColor(255, 255, 255)    # White
ACCENT_COLOR = RGBColor(0, 112, 192)          # Blue

# Font Configuration
TITLE_FONT = "Calibri"
CONTENT_FONT = "Calibri"

# Font Sizes
TITLE_SLIDE_TITLE_SIZE = Pt(54)
TITLE_SLIDE_SUBTITLE_SIZE = Pt(28)
CONTENT_SLIDE_TITLE_SIZE = Pt(36)
BULLET_LEVEL_1_SIZE = Pt(20)
BULLET_LEVEL_2_SIZE = Pt(18)

# Layout and Spacing
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)
TITLE_TOP = Inches(2.0)
TITLE_LEFT = Inches(0.5)
TITLE_WIDTH = Inches(9)
TITLE_HEIGHT = Inches(1.5)
CONTENT_TOP = Inches(1.2)
CONTENT_LEFT = Inches(0.5)
CONTENT_WIDTH = Inches(9)
CONTENT_HEIGHT = Inches(5.5)
BULLET_INDENT = Inches(0.5)

# Text Limits
MAX_TITLE_LENGTH = 100
MAX_SUBTITLE_LENGTH = 150
MAX_BULLET_LENGTH = 200


# ==================== HELPER FUNCTIONS ====================

def generate_unique_filename(base_name=FILENAME_PREFIX):
    """
    Generate a unique filename for the presentation
    
    Args:
        base_name (str): Base name for the file
        
    Returns:
        str: Unique filename with timestamp and UUID
    """
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    unique_id = str(uuid.uuid4())[:8]
    filename = f"{base_name}_{timestamp}_{unique_id}.pptx"
    return filename


def truncate_text(text, max_length):
    """
    Truncate text to maximum length with ellipsis
    
    Args:
        text (str): Text to truncate
        max_length (int): Maximum length
        
    Returns:
        str: Truncated text
    """
    text = str(text).strip()
    if len(text) > max_length:
        return text[:max_length-3] + "..."
    return text


def sanitize_text(text):
    """
    Clean and sanitize text for presentation
    
    Args:
        text (str): Input text
        
    Returns:
        str: Sanitized text
    """
    if not text:
        return ""
    
    text = str(text).strip()
    # Remove excessive whitespace
    text = ' '.join(text.split())
    return text


def ensure_output_folder(folder_path):
    """
    Ensure output folder exists
    
    Args:
        folder_path (str): Path to output folder
    """
    os.makedirs(folder_path, exist_ok=True)


# ==================== SLIDE CREATION FUNCTIONS ====================

def add_title_slide(prs, slide_data):
    """
    Add a title slide to the presentation
    
    Args:
        prs: Presentation object
        slide_data (dict): Data for the slide with 'title' and 'subtitle'
        
    Returns:
        Slide object
    """
    # Use blank layout to have more control
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Get and sanitize text
    title_text = sanitize_text(slide_data.get('title', 'Untitled Presentation'))
    subtitle_text = sanitize_text(slide_data.get('subtitle', ''))
    
    # Truncate if necessary
    title_text = truncate_text(title_text, MAX_TITLE_LENGTH)
    subtitle_text = truncate_text(subtitle_text, MAX_SUBTITLE_LENGTH)
    
    # Add title text box
    title_left = Inches(0.5)
    title_top = Inches(2.5)
    title_width = Inches(9)
    title_height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.word_wrap = True
    
    # Format title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.font.name = TITLE_FONT
    title_paragraph.font.size = TITLE_SLIDE_TITLE_SIZE
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = TITLE_COLOR
    
    # Add subtitle if exists
    if subtitle_text:
        subtitle_left = Inches(0.5)
        subtitle_top = Inches(4.2)
        subtitle_width = Inches(9)
        subtitle_height = Inches(1.0)
        
        subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_frame.word_wrap = True
        
        # Format subtitle
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        subtitle_paragraph.font.name = TITLE_FONT
        subtitle_paragraph.font.size = TITLE_SLIDE_SUBTITLE_SIZE
        subtitle_paragraph.font.color.rgb = SUBTITLE_COLOR
    
    return slide


def add_content_slide(prs, slide_data):
    """
    Add a content slide with title and bullet points
    
    Args:
        prs: Presentation object
        slide_data (dict): Data with 'title' and 'bullets'
        
    Returns:
        Slide object
    """
    # Use title and content layout
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Get and sanitize title
    title_text = sanitize_text(slide_data.get('title', 'Slide Title'))
    title_text = truncate_text(title_text, MAX_TITLE_LENGTH)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title_text
    
    # Format title
    title_frame = title_shape.text_frame
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = CONTENT_FONT
    title_paragraph.font.size = CONTENT_SLIDE_TITLE_SIZE
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = CONTENT_TITLE_COLOR
    
    # Get bullets
    bullets = slide_data.get('bullets', [])
    
    if bullets:
        # Get content placeholder
        content_shape = slide.shapes.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()  # Clear default text
        text_frame.word_wrap = True
        
        # Add each bullet
        for i, bullet_text in enumerate(bullets):
            bullet_text = sanitize_text(bullet_text)
            bullet_text = truncate_text(bullet_text, MAX_BULLET_LENGTH)
            
            if not bullet_text:
                continue
            
            # Add paragraph
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            paragraph.text = bullet_text
            paragraph.level = 0  # First level bullet
            paragraph.font.name = CONTENT_FONT
            paragraph.font.size = BULLET_LEVEL_1_SIZE
            paragraph.font.color.rgb = BULLET_COLOR
            paragraph.space_after = Pt(12)
    
    return slide


def apply_presentation_settings(prs):
    """
    Apply global settings to the presentation
    
    Args:
        prs: Presentation object
    """
    # Set slide size (standard 16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)


# ==================== MAIN GENERATION FUNCTION ====================

def create_presentation(structure, output_folder=OUTPUT_FOLDER):
    """
    Create a PowerPoint presentation from structured data
    
    Args:
        structure (dict): Structured presentation data with title and slides
        output_folder (str): Folder to save the presentation
        
    Returns:
        str: Full path to the generated presentation file
        
    Raises:
        Exception: If presentation cannot be created
    """
    try:
        # Validate input structure
        if not structure or not isinstance(structure, dict):
            raise ValueError("Invalid presentation structure")
        
        if 'slides' not in structure or not structure['slides']:
            raise ValueError("No slides found in structure")
        
        # Ensure output folder exists
        ensure_output_folder(output_folder)
        
        # Create new presentation
        prs = Presentation()
        apply_presentation_settings(prs)
        
        # Process each slide
        slides_added = 0
        for slide_data in structure['slides']:
            slide_type = slide_data.get('type', 'content')
            
            try:
                if slide_type == 'title':
                    add_title_slide(prs, slide_data)
                    slides_added += 1
                elif slide_type == 'content':
                    add_content_slide(prs, slide_data)
                    slides_added += 1
                else:
                    print(f"Warning: Unknown slide type '{slide_type}', skipping")
            except Exception as e:
                print(f"Warning: Failed to add slide: {str(e)}")
                continue
        
        if slides_added == 0:
            raise Exception("No slides were successfully added to presentation")
        
        # Generate filename
        filename = generate_unique_filename()
        file_path = os.path.join(output_folder, filename)
        
        # Save presentation
        prs.save(file_path)
        
        print(f"✓ Presentation created successfully: {filename}")
        print(f"  Total slides: {slides_added}")
        print(f"  File size: {os.path.getsize(file_path) / 1024:.1f} KB")
        
        return file_path
    
    except Exception as e:
        raise Exception(f"Failed to create presentation: {str(e)}")


def validate_structure(structure):
    """
    Validate presentation structure before generation
    
    Args:
        structure (dict): Presentation structure to validate
        
    Returns:
        tuple: (is_valid: bool, error_message: str or None)
    """
    if not structure or not isinstance(structure, dict):
        return False, "Structure must be a dictionary"
    
    if 'title' not in structure:
        return False, "Missing 'title' field"
    
    if 'slides' not in structure:
        return False, "Missing 'slides' field"
    
    if not isinstance(structure['slides'], list):
        return False, "'slides' must be a list"
    
    if len(structure['slides']) == 0:
        return False, "No slides in structure"
    
    # Validate each slide
    for i, slide in enumerate(structure['slides']):
        if not isinstance(slide, dict):
            return False, f"Slide {i+1} must be a dictionary"
        
        if 'type' not in slide:
            return False, f"Slide {i+1} missing 'type' field"
        
        if 'title' not in slide:
            return False, f"Slide {i+1} missing 'title' field"
        
        if slide['type'] == 'content' and 'bullets' not in slide:
            return False, f"Content slide {i+1} missing 'bullets' field"
    
    return True, None


def get_presentation_info(file_path):
    """
    Get information about a generated presentation
    
    Args:
        file_path (str): Path to presentation file
        
    Returns:
        dict: Information about the presentation
    """
    if not os.path.exists(file_path):
        return None
    
    try:
        prs = Presentation(file_path)
        
        return {
            'filename': os.path.basename(file_path),
            'file_path': file_path,
            'file_size_kb': os.path.getsize(file_path) / 1024,
            'slide_count': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height
        }
    except Exception as e:
        return {'error': str(e)}


def cleanup_old_presentations(folder_path, max_age_hours=24):
    """
    Clean up old presentation files
    
    Args:
        folder_path (str): Folder containing presentations
        max_age_hours (int): Maximum age in hours before deletion
    """
    import time
    
    if not os.path.exists(folder_path):
        return
    
    current_time = time.time()
    max_age_seconds = max_age_hours * 3600
    
    for filename in os.listdir(folder_path):
        if filename.endswith('.pptx'):
            file_path = os.path.join(folder_path, filename)
            file_age = current_time - os.path.getmtime(file_path)
            
            if file_age > max_age_seconds:
                try:
                    os.remove(file_path)
                    print(f"Cleaned up old file: {filename}")
                except Exception as e:
                    print(f"Failed to delete {filename}: {str(e)}")


def create_presentation_with_template(structure, template_path, output_folder="temp"):
    """
    Create a PowerPoint presentation using a template file
    
    Args:
        structure (dict): Structured presentation data with title and slides
        template_path (str): Path to the template PPTX file
        output_folder (str): Folder to save the presentation
        
    Returns:
        str: Full path to the generated presentation file
        
    Raises:
        Exception: If presentation cannot be created
    """
    try:
        # Validate template exists
        if not os.path.exists(template_path):
            print(f"Warning: Template not found at {template_path}, falling back to default generation")
            return create_presentation(structure, output_folder)
        
        # Validate input structure
        if not structure or not isinstance(structure, dict):
            raise ValueError("Invalid presentation structure")
        
        if 'slides' not in structure or not structure['slides']:
            raise ValueError("No slides found in structure")
        
        # Ensure output folder exists
        os.makedirs(output_folder, exist_ok=True)
        
        print(f"Creating presentation with template: {template_path}")
        
        # Load the template
        prs = Presentation(template_path)
        
        # Get slide layouts from template
        slide_layouts = prs.slide_layouts
        
        # Determine which layouts to use
        title_layout = slide_layouts[0] if len(slide_layouts) > 0 else None
        content_layout = slide_layouts[1] if len(slide_layouts) > 1 else slide_layouts[0]
        section_layout = slide_layouts[2] if len(slide_layouts) > 2 else content_layout
        
        # Clear existing slides from template (keep only the layouts)
        while len(prs.slides) > 0:
            xml_slides = prs.slides._sldIdLst
            slides_list = list(xml_slides)
            xml_slides.remove(slides_list[0])
        
        # Track slides added
        slides_added = 0
        
        # Process each slide in the structure
        for slide_data in structure['slides']:
            slide_type = slide_data.get('type', 'content')
            
            try:
                if slide_type == 'title':
                    # Add title slide
                    slide = prs.slides.add_slide(title_layout)
                    _populate_title_slide(slide, slide_data)
                    slides_added += 1
                    
                elif slide_type == 'section':
                    # Add section divider slide
                    slide = prs.slides.add_slide(section_layout)
                    _populate_section_slide(slide, slide_data)
                    slides_added += 1
                    
                elif slide_type == 'content':
                    # Add content slide with bullets
                    slide = prs.slides.add_slide(content_layout)
                    _populate_content_slide(slide, slide_data)
                    slides_added += 1
                    
                elif slide_type == 'two-column':
                    # Add two-column slide if template supports it
                    two_col_layout = slide_layouts[3] if len(slide_layouts) > 3 else content_layout
                    slide = prs.slides.add_slide(two_col_layout)
                    _populate_two_column_slide(slide, slide_data)
                    slides_added += 1
                    
                elif slide_type == 'image':
                    # Add image slide if template supports it
                    image_layout = slide_layouts[4] if len(slide_layouts) > 4 else content_layout
                    slide = prs.slides.add_slide(image_layout)
                    _populate_image_slide(slide, slide_data)
                    slides_added += 1
                    
                else:
                    # Default to content slide for unknown types
                    slide = prs.slides.add_slide(content_layout)
                    _populate_content_slide(slide, slide_data)
                    slides_added += 1
                    
            except Exception as e:
                print(f"Warning: Failed to add slide: {str(e)}")
                continue
        
        if slides_added == 0:
            raise Exception("No slides were successfully added to presentation")
        
        # Generate filename
        title = structure.get('title', 'Presentation')
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = generate_unique_filename(safe_title)
        file_path = os.path.join(output_folder, filename)
        
        # Save presentation
        prs.save(file_path)
        
        print(f"✓ Presentation created with template successfully: {filename}")
        print(f"  Total slides: {slides_added}")
        print(f"  File size: {os.path.getsize(file_path) / 1024:.1f} KB")
        
        return file_path
        
    except Exception as e:
        print(f"Error creating presentation with template: {str(e)}")
        print("Falling back to default presentation generation")
        return create_presentation(structure, output_folder)


def _populate_title_slide(slide, slide_data):
    """
    Populate a title slide with data
    
    Args:
        slide: PowerPoint slide object
        slide_data (dict): Data containing title and subtitle
    """
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get('title', 'Untitled')
    
    # Set subtitle if available
    if 'subtitle' in slide_data:
        # Find subtitle placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle is usually index 1
                shape.text = slide_data['subtitle']
                break


def _populate_section_slide(slide, slide_data):
    """
    Populate a section divider slide
    
    Args:
        slide: PowerPoint slide object
        slide_data (dict): Data containing section name
    """
    # Set section title
    if slide.shapes.title:
        section_name = slide_data.get('name', slide_data.get('title', 'Section'))
        slide.shapes.title.text = section_name
    
    # Add section description if available
    if 'description' in slide_data:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = slide_data['description']
                break


def _populate_content_slide(slide, slide_data):
    """
    Populate a content slide with bullets
    
    Args:
        slide: PowerPoint slide object
        slide_data (dict): Data containing title and bullets
    """
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get('title', 'Slide Title')
    
    # Find content placeholder and add bullets
    for shape in slide.placeholders:
        if shape.has_text_frame and shape.placeholder_format.idx == 1:
            text_frame = shape.text_frame
            text_frame.clear()  # Clear any existing text
            
            bullets = slide_data.get('bullets', [])
            if not bullets and 'content' in slide_data:
                # If no bullets but content exists, split content into bullets
                content = slide_data['content']
                bullets = content.split('\n') if isinstance(content, str) else [str(content)]
            
            # Add bullets
            for i, bullet in enumerate(bullets):
                if i == 0:
                    # Use first paragraph
                    p = text_frame.paragraphs[0]
                else:
                    # Add new paragraph
                    p = text_frame.add_paragraph()
                
                p.text = str(bullet).strip()
                p.level = 0  # Main bullet level
            
            break
    
    # Add speaker notes if available
    if 'notes' in slide_data and slide_data['notes']:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data['notes']


def _populate_two_column_slide(slide, slide_data):
    """
    Populate a two-column slide
    
    Args:
        slide: PowerPoint slide object
        slide_data (dict): Data containing title, left_content, and right_content
    """
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get('title', 'Slide Title')
    
    # Find and populate left and right content placeholders
    placeholders = list(slide.placeholders)
    
    # Left column (usually index 1)
    if len(placeholders) > 1 and 'left_content' in slide_data:
        left_placeholder = placeholders[1]
        if left_placeholder.has_text_frame:
            _add_bullets_to_placeholder(left_placeholder, slide_data['left_content'])
    
    # Right column (usually index 2)
    if len(placeholders) > 2 and 'right_content' in slide_data:
        right_placeholder = placeholders[2]
        if right_placeholder.has_text_frame:
            _add_bullets_to_placeholder(right_placeholder, slide_data['right_content'])


def _populate_image_slide(slide, slide_data):
    """
    Populate an image slide
    
    Args:
        slide: PowerPoint slide object
        slide_data (dict): Data containing title and image_path
    """
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get('title', 'Image')
    
    # Add caption if available
    if 'caption' in slide_data:
        for shape in slide.placeholders:
            if shape.has_text_frame and shape.placeholder_format.idx > 0:
                shape.text = slide_data['caption']
                break
    
    # Note: Actual image insertion would require the image file path
    # This is a placeholder for future implementation
    if 'image_path' in slide_data and os.path.exists(slide_data['image_path']):
        # TODO: Implement image insertion
        pass


def _add_bullets_to_placeholder(placeholder, content):
    """
    Helper function to add bullets to a placeholder
    
    Args:
        placeholder: PowerPoint placeholder shape
        content: List of bullets or string content
    """
    if not placeholder.has_text_frame:
        return
    
    text_frame = placeholder.text_frame
    text_frame.clear()
    
    # Convert content to list if string
    if isinstance(content, str):
        bullets = content.split('\n')
    elif isinstance(content, list):
        bullets = content
    else:
        bullets = [str(content)]
    
    # Add bullets
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = str(bullet).strip()
        p.level = 0


def apply_template_to_existing(presentation_path, template_path):
    """
    Apply a template to an existing presentation
    Note: This is complex and may not preserve all content perfectly
    
    Args:
        presentation_path (str): Path to existing presentation
        template_path (str): Path to template file
        
    Returns:
        str: Path to the new presentation with template applied
    """
    try:
        if not os.path.exists(presentation_path):
            raise FileNotFoundError(f"Presentation not found: {presentation_path}")
        
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
        
        # Load both presentations
        prs_content = Presentation(presentation_path)
        prs_template = Presentation(template_path)
        
        # Create new presentation from template
        prs_new = Presentation(template_path)
        
        # Clear template slides
        while len(prs_new.slides) > 0:
            xml_slides = prs_new.slides._sldIdLst
            slides_list = list(xml_slides)
            xml_slides.remove(slides_list[0])
        
        # Get layouts from template
        layouts = prs_new.slide_layouts
        
        # Copy content from original to new presentation with template layouts
        for slide in prs_content.slides:
            # Determine best layout to use
            if slide.slide_layout.name == 'Title Slide':
                new_layout = layouts[0]
            elif 'Title and Content' in slide.slide_layout.name:
                new_layout = layouts[1] if len(layouts) > 1 else layouts[0]
            elif 'Section' in slide.slide_layout.name:
                new_layout = layouts[2] if len(layouts) > 2 else layouts[1]
            else:
                new_layout = layouts[1] if len(layouts) > 1 else layouts[0]
            
            # Add new slide with template layout
            new_slide = prs_new.slides.add_slide(new_layout)
            
            # Copy title
            if slide.shapes.title and new_slide.shapes.title:
                new_slide.shapes.title.text = slide.shapes.title.text
            
            # Copy content (simplified - may need enhancement for complex slides)
            for shape in slide.shapes:
                if shape.has_text_frame and not shape.shape_id == slide.shapes.title.shape_id:
                    # Find corresponding placeholder in new slide
                    for new_shape in new_slide.shapes:
                        if new_shape.has_text_frame and not new_shape.shape_id == new_slide.shapes.title.shape_id:
                            # Copy text content
                            new_shape.text_frame.clear()
                            for paragraph in shape.text_frame.paragraphs:
                                new_p = new_shape.text_frame.add_paragraph()
                                new_p.text = paragraph.text
                                new_p.level = paragraph.level
                            break
        
        # Generate output filename
        base_name = os.path.splitext(os.path.basename(presentation_path))[0]
        output_path = presentation_path.replace('.pptx', '_templated.pptx')
        
        # Save new presentation
        prs_new.save(output_path)
        
        print(f"✓ Template applied successfully: {output_path}")
        return output_path
        
    except Exception as e:
        print(f"Error applying template: {str(e)}")
        return presentation_path  # Return original if template application fails